#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Design-Build Phase Research
AI Code Assistants for Federal Environments
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

# Color scheme
BRAND_PURPLE = RGBColor(117, 0, 192)  # #7500c0
ACCENT_PURPLE = RGBColor(160, 85, 245)  # #a055f5
DARK_GRAY = RGBColor(51, 51, 51)  # #333333
LIGHT_GRAY = RGBColor(128, 128, 128)  # #808080
WHITE = RGBColor(255, 255, 255)
PROCEED_GREEN = RGBColor(16, 185, 129)  # #10B981
CAUTION_AMBER = RGBColor(245, 158, 11)  # #F59E0B


def set_slide_background(slide, color):
    """Set solid background color for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle):
    """Add title slide with gradient-style header"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Purple header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(2.5)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BRAND_PURPLE
    header.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.8),
        Inches(9), Inches(1.2)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.9),
        Inches(9), Inches(0.5)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    # Footer
    footer_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.1),
        Inches(9), Inches(0.3)
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GenAI COTS Team | Accenture Federal Services | December 2025"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY

    return slide


def add_section_header(prs, title, subtitle=None):
    """Add section header slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Purple accent bar on left
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.15), Inches(5.63)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND_PURPLE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2),
        Inches(9), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BRAND_PURPLE

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY

    return slide


def add_content_slide(prs, title, content_items, highlight_first=False):
    """Add content slide with bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BRAND_PURPLE
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25),
        Inches(9), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2),
        Inches(9), Inches(4)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)

        if highlight_first and i == 0:
            p.font.bold = True
            p.font.color.rgb = BRAND_PURPLE

    return slide


def add_table_slide(prs, title, headers, rows, highlight_col=None):
    """Add slide with table"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BRAND_PURPLE
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.15),
        Inches(9.4), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Calculate table dimensions
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)

    table_width = Inches(9.4)
    table_height = Inches(4)

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.3), Inches(1),
        table_width, table_height
    ).table

    # Style header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_PURPLE

        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER

        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Fill data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_value)

            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 245, 250)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

            # Highlight column if specified
            if highlight_col is not None and col_idx == highlight_col:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(243, 232, 255)  # Light purple

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = DARK_GRAY
                paragraph.alignment = PP_ALIGN.CENTER

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return slide


def add_tool_highlight_slide(prs, tool_name, score, path, features, recommendation):
    """Add individual tool highlight slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BRAND_PURPLE
    title_bar.line.fill.background()

    # Tool name
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(6), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = tool_name
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Score badge
    score_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.5), Inches(0.15),
        Inches(2.2), Inches(0.7)
    )
    score_box.fill.solid()
    score_box.fill.fore_color.rgb = PROCEED_GREEN if float(score.split('/')[0]) >= 8.5 else CAUTION_AMBER
    score_box.line.fill.background()

    score_text = slide.shapes.add_textbox(
        Inches(7.5), Inches(0.25),
        Inches(2.2), Inches(0.5)
    )
    tf = score_text.text_frame
    p = tf.paragraphs[0]
    p.text = f"Score: {score}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Path indicator
    path_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1),
        Inches(9), Inches(0.4)
    )
    tf = path_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Path: {path}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY

    # Key features
    features_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.6),
        Inches(9), Inches(2.5)
    )
    tf = features_box.text_frame
    tf.word_wrap = True

    for i, feature in enumerate(features):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✓ {feature}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

    # Recommendation box
    rec_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(4.3),
        Inches(9), Inches(0.8)
    )
    rec_box.fill.solid()
    rec_box.fill.fore_color.rgb = RGBColor(243, 232, 255)  # Light purple
    rec_box.line.color.rgb = BRAND_PURPLE

    rec_text = slide.shapes.add_textbox(
        Inches(0.7), Inches(4.45),
        Inches(8.6), Inches(0.5)
    )
    tf = rec_text.text_frame
    p = tf.paragraphs[0]
    p.text = f"Recommendation: {recommendation}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BRAND_PURPLE

    return slide


def add_decision_matrix_slide(prs):
    """Add when-to-choose decision matrix"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BRAND_PURPLE
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.15),
        Inches(9.4), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Decision Matrix: When to Choose Each Tool"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Decision items
    decisions = [
        ("Maximum air-gap security", "Tabby", "Zero cloud deps, Rust, local models"),
        ("FedRAMP SaaS with full features", "Windsurf", "Only FedRAMP High certified"),
        ("Open-source flexibility", "Continue.dev", "Apache 2.0, 20+ LLM providers"),
        ("Terminal/CLI workflow", "Claude Code", "Native terminal, agentic"),
        ("Enterprise air-gap + support", "Tabnine", "SOC 2, IP indemnification"),
        ("Code review + testing focus", "Qodo", "PR-Agent, test generation"),
    ]

    y_pos = 1.1
    for need, choice, reason in decisions:
        # Need
        need_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(y_pos),
            Inches(3.2), Inches(0.5)
        )
        tf = need_box.text_frame
        p = tf.paragraphs[0]
        p.text = need
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY

        # Arrow
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(3.5), Inches(y_pos + 0.1),
            Inches(0.4), Inches(0.25)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = BRAND_PURPLE
        arrow.line.fill.background()

        # Choice
        choice_box = slide.shapes.add_textbox(
            Inches(4), Inches(y_pos),
            Inches(1.8), Inches(0.5)
        )
        tf = choice_box.text_frame
        p = tf.paragraphs[0]
        p.text = choice
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BRAND_PURPLE

        # Reason
        reason_box = slide.shapes.add_textbox(
            Inches(5.8), Inches(y_pos),
            Inches(4), Inches(0.5)
        )
        tf = reason_box.text_frame
        p = tf.paragraphs[0]
        p.text = reason
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = LIGHT_GRAY

        y_pos += 0.65

    return slide


def add_summary_slide(prs):
    """Add final summary slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Full purple background
    set_slide_background(slide, BRAND_PURPLE)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Takeaways"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    takeaways = [
        "Air-Gapped/Classified: Tabby is the gold standard",
        "AWS GovCloud: Claude Code via Bedrock offers FedRAMP High",
        "Flexibility: Continue.dev provides maximum LLM choice",
        "Enterprise SaaS: Windsurf is the only FedRAMP High certified",
        "All six tools provide viable federal paths"
    ]

    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.8),
        Inches(9), Inches(3)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, takeaway in enumerate(takeaways):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"→ {takeaway}"
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(16)

    # Footer
    footer_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.1),
        Inches(9), Inches(0.3)
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Contact: christopher.g.roge@afs.com"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 200, 220)

    return slide


def main():
    """Generate the Design-Build research presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)  # 16:9 aspect ratio

    # 1. Title Slide
    add_title_slide(
        prs,
        "AI Code Assistants for Federal Environments",
        "Design-Build Phase | Agentic SDLC Research"
    )

    # 2. Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Evaluated 22 AI code assistants for federal deployment viability",
        "Two deployment paths: Path A (FedRAMP SaaS) and Path B (Self-Hosted with GovCloud LLMs)",
        "Top recommendations: Tabby (9/10), Continue.dev (8.5/10), Claude Code (8.5/10), Windsurf (8.5/10)",
        "All six evaluated tools provide viable federal paths",
        "Choice depends on security requirements, infrastructure, and feature priorities"
    ], highlight_first=True)

    # 3. Federal Deployment Paths
    add_section_header(prs, "Federal Deployment Paths", "Understanding Path A vs Path B")

    add_content_slide(prs, "Path A: FedRAMP Cloud Deployment", [
        "Tool vendor provides FedRAMP-authorized cloud service",
        "Lower infrastructure burden - vendor manages compute",
        "Limited to unclassified workloads (IL2-IL4)",
        "Currently only Windsurf has actual FedRAMP High authorization",
        "GitHub Copilot expected to follow (verify status)"
    ])

    add_content_slide(prs, "Path B: Self-Hosted with GovCloud LLMs", [
        "Tool runs on-premises or in agency cloud",
        "Connects to authorized LLM backends (AWS Bedrock, Azure OpenAI)",
        "Supports air-gapped and classified workloads (IL5+)",
        "Higher infrastructure responsibility",
        "Multiple excellent options: Tabby, Continue.dev, Tabnine, Qodo"
    ])

    # 4. Federal Viability Matrix
    add_section_header(prs, "Federal Viability Matrix", "Comparing Top 6 Candidates")

    add_table_slide(
        prs,
        "Federal Compliance Comparison",
        ["Tool", "Score", "Path", "FedRAMP", "Air-Gap", "GovCloud LLM"],
        [
            ["Tabby", "9/10", "B", "N/A", "Yes", "Bedrock, Azure"],
            ["Continue.dev", "8.5/10", "B", "N/A", "Yes", "Bedrock, Azure"],
            ["Claude Code", "8.5/10", "B", "Via Bedrock", "No", "Bedrock"],
            ["Windsurf", "8.5/10", "A+B", "High", "BYOL only", "Bedrock, Azure"],
            ["Tabnine", "8/10", "B", "No", "Yes", "Bedrock, Azure"],
            ["Qodo", "8/10", "B", "No", "Yes", "Bedrock, Azure"],
        ],
        highlight_col=1
    )

    add_table_slide(
        prs,
        "Architecture & Capabilities",
        ["Tool", "Architecture", "Self-Hosted", "Agentic", "Code Review"],
        [
            ["Tabby", "Rust binary", "Full features", "Limited", "No"],
            ["Continue.dev", "IDE extension", "Full features", "Yes", "No"],
            ["Claude Code", "CLI client", "N/A", "Native", "Yes"],
            ["Windsurf", "IDE (fork)", "Reduced*", "Cascade", "No"],
            ["Tabnine", "IDE extension", "Full features", "Limited", "No"],
            ["Qodo", "IDE extension", "Full features", "Yes", "PR-Agent"],
        ]
    )

    # 5. Tool Highlights
    add_section_header(prs, "Top Recommendations", "Detailed Tool Profiles")

    add_tool_highlight_slide(
        prs,
        "Tabby",
        "9/10",
        "B (Self-Hosted)",
        [
            "Zero cloud dependencies - fully air-gapped capable",
            "Written in Rust for memory safety and performance",
            "Apache 2.0 open-source - no vendor lock-in",
            "Supports AWS Bedrock via OpenAI-compatible API",
            "~$500-2K/month infrastructure for 50 engineers"
        ],
        "PROCEED - Best for air-gapped/IL5+ environments"
    )

    add_tool_highlight_slide(
        prs,
        "Continue.dev",
        "8.5/10",
        "B (Self-Hosted)",
        [
            "Apache 2.0 open-source with 30K+ GitHub stars",
            "Native AWS Bedrock and Azure OpenAI providers",
            "Full offline mode with Ollama local models",
            "20+ LLM providers supported",
            "Enterprise: SSO, governance, managed proxy"
        ],
        "PROCEED - Best open-source option with GovCloud flexibility"
    )

    add_tool_highlight_slide(
        prs,
        "Claude Code",
        "8.5/10",
        "B (Via Bedrock)",
        [
            "CLAUDE_CODE_USE_BEDROCK=1 routes to GovCloud",
            "FedRAMP High, DoD IL4/5 via Bedrock",
            "Telemetry auto-disabled with Bedrock",
            "Terminal-native (no IDE required)",
            "MCP protocol for enterprise integrations"
        ],
        "CONDITIONAL - Best for AWS GovCloud users (requires internet)"
    )

    add_tool_highlight_slide(
        prs,
        "Windsurf",
        "8.5/10",
        "A+B (Hybrid)",
        [
            "FedRAMP High authorized (March 2025) via Palantir FedStart",
            "DoD IL4, IL5, IL6, ITAR compliant",
            "Hybrid deployment with full features",
            "Self-hosted loses Cascade agentic feature",
            "$60/user/month, 2-4 week setup"
        ],
        "PROCEED - Only actual FedRAMP High certified option"
    )

    # 6. Path Recommendations
    add_section_header(prs, "Path-Based Recommendations", "Match Tools to Requirements")

    add_table_slide(
        prs,
        "Recommendations by Environment",
        ["Environment", "Rank 1", "Rank 2", "Rank 3"],
        [
            ["Air-Gapped/IL5+", "Tabby", "Tabnine", "Continue+Ollama"],
            ["AWS GovCloud", "Claude Code", "Continue.dev", "Windsurf"],
            ["Azure GovCloud", "Continue.dev", "Tabnine", "Qodo"],
            ["FedRAMP SaaS", "Windsurf", "GitHub Copilot*", "Amazon Q*"],
        ]
    )

    # 7. Decision Matrix
    add_decision_matrix_slide(prs)

    # 8. Next Steps
    add_content_slide(prs, "Recommended Next Steps", [
        "POC Deployment: Deploy Tabby or Continue.dev in non-production environment (1-3 days)",
        "Pilot Program: 5-10 developers for 2-4 weeks with metrics collection",
        "Security Review: Source code audit for classified deployments",
        "Infrastructure Planning: GPU provisioning for self-hosted options",
        "Vendor Engagement: Contact Windsurf for FedRAMP hybrid deployment timeline"
    ])

    # 9. Summary
    add_summary_slide(prs)

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "00-DESIGN-BUILD-RESEARCH.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()

#!/usr/bin/env python3
"""
Create PowerPoint slides from 01-one-pager.html files
Faithfully recreates the HTML design in a single slide PowerPoint
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import re
from pathlib import Path


def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def remove_shadow(shape):
    """Completely remove shadow from a shape by clearing XML elements"""
    try:
        # Access the shape's XML element
        sp = shape._element
        spPr = sp.spPr

        # Find and remove effect list elements (which contain shadows)
        for elem in list(spPr):
            if 'effectLst' in str(elem.tag) or 'outerShdw' in str(elem.tag):
                spPr.remove(elem)
    except:
        pass  # If XML manipulation fails, continue anyway


def parse_html_one_pager(html_path):
    """Parse HTML one-pager and extract content"""
    with open(html_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')

    data = {
        'title': soup.find('h1').text if soup.find('h1') else '',
        'overview': '',
        'overview_highlights': [],
        'callouts': [],
        'table_headers': [],
        'table_content': [],
        'footer_contact': ''
    }

    # Extract overview with highlights
    overview_div = soup.find('div', class_='overview-text')
    if overview_div:
        data['overview'] = overview_div.get_text()
        highlights = overview_div.find_all('span', class_='purple-highlight')
        data['overview_highlights'] = [h.text for h in highlights]

    # Extract callouts
    callouts = soup.find_all('div', class_='callout-container')
    for callout in callouts:
        title_elem = callout.find('div', class_='callout-title')
        content_elem = callout.find('div', class_='callout-content')

        # Get callout type from class
        callout_type = 'default'
        if 'critical-finding' in callout.get('class', []):
            callout_type = 'critical'
        elif 'recommendation' in callout.get('class', []):
            callout_type = 'recommendation'
        elif 'bottom-line' in callout.get('class', []):
            callout_type = 'bottom-line'

        if title_elem and content_elem:
            # Remove SVG text from title
            title_text = title_elem.get_text(strip=True)
            data['callouts'].append({
                'type': callout_type,
                'title': title_text,
                'content': content_elem.get_text(strip=True)
            })

    # Extract table
    table = soup.find('table')
    if table:
        headers = table.find_all('th')
        data['table_headers'] = [h.text for h in headers]

        rows = table.find_all('tr')
        for row in rows[1:]:  # Skip header row
            cells = row.find_all('td')
            data['table_content'].append([c.get_text(strip=True) for c in cells])

    # Extract footer contact
    footer_info = soup.find('div', class_='footer-info')
    if footer_info:
        data['footer_contact'] = footer_info.get_text(separator=' | ', strip=True)

    return data


def create_pptx_slide(data, output_path, tool_name):
    """Create PowerPoint slide from parsed data"""

    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # Add blank slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Color definitions
    PURPLE_CORE = RGBColor(117, 0, 192)      # #7500c0
    PURPLE_ACCENT = RGBColor(160, 85, 245)   # #a055f5
    BLACK = RGBColor(0, 0, 0)
    GRAY = RGBColor(127, 140, 141)
    AMBER = RGBColor(217, 119, 6)            # #D97706
    GREEN = RGBColor(16, 185, 129)           # #10B981

    # Gradient colors for callouts
    GRADIENT_AMBER_START = RGBColor(254, 243, 199)  # #FEF3C7
    GRADIENT_GREEN_START = RGBColor(209, 250, 229)  # #D1FAE5
    GRADIENT_PURPLE_START = RGBColor(243, 232, 255) # #F3E8FF

    # Layout margins
    MARGIN = Inches(0.5)
    CONTENT_WIDTH = prs.slide_width - (2 * MARGIN)

    y_pos = MARGIN

    # Title
    title_box = slide.shapes.add_textbox(
        MARGIN, y_pos,
        CONTENT_WIDTH, Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = data['title']
    p.font.name = 'Graphik'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # Add underline manually with a shape
    underline = slide.shapes.add_shape(
        1,  # Rectangle
        MARGIN, y_pos + Inches(0.55),
        CONTENT_WIDTH, Inches(0.03)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = PURPLE_CORE
    underline.line.fill.background()  # No outline

    y_pos += Inches(0.75)

    # Overview text
    overview_box = slide.shapes.add_textbox(
        MARGIN, y_pos,
        CONTENT_WIDTH, Inches(0.8)
    )
    overview_frame = overview_box.text_frame
    overview_frame.word_wrap = True
    p = overview_frame.paragraphs[0]

    # Add overview with purple highlights
    overview_text = data['overview']
    for highlight in data['overview_highlights']:
        overview_text = overview_text.replace(highlight, f'[[{highlight}]]')

    parts = re.split(r'\[\[(.*?)\]\]', overview_text)
    for i, part in enumerate(parts):
        if part in data['overview_highlights']:
            run = p.add_run()
            run.text = part
            run.font.name = 'Graphik'
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = PURPLE_CORE
        elif part:
            run = p.add_run()
            run.text = part
            run.font.name = 'Graphik'
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = BLACK

    y_pos += Inches(1.0)

    # Callouts (first 2: critical finding and recommendation)
    callout_height = Inches(0.7)
    for callout in data['callouts'][:2]:
        # Determine colors
        if callout['type'] == 'critical':
            bg_color = GRADIENT_AMBER_START
            border_color = AMBER
        elif callout['type'] == 'recommendation':
            bg_color = GRADIENT_GREEN_START
            border_color = GREEN
        else:
            bg_color = GRADIENT_PURPLE_START
            border_color = PURPLE_CORE

        # Background box
        bg_box = slide.shapes.add_shape(
            1,  # Rectangle
            MARGIN, y_pos,
            CONTENT_WIDTH, callout_height
        )
        bg_box.fill.solid()
        bg_box.fill.fore_color.rgb = bg_color
        bg_box.line.fill.background()  # No outline
        remove_shadow(bg_box)  # Remove shadow completely

        # Title (left side, bold)
        title_width = Inches(2.5)
        title_box = slide.shapes.add_textbox(
            MARGIN + Inches(0.2), y_pos + Inches(0.15),
            title_width, Inches(0.4)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = callout['title']
        p.font.name = 'Graphik'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLACK

        # Content (right side)
        content_box = slide.shapes.add_textbox(
            MARGIN + title_width + Inches(0.3), y_pos + Inches(0.1),
            CONTENT_WIDTH - title_width - Inches(0.5), callout_height - Inches(0.2)
        )
        cf = content_box.text_frame
        cf.word_wrap = True
        cf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = cf.paragraphs[0]
        cp.text = callout['content']
        cp.font.name = 'Graphik'
        cp.font.size = Pt(12)
        cp.font.color.rgb = BLACK

        y_pos += callout_height + Inches(0.15)

    # Table
    if data['table_headers'] and data['table_content']:
        table_top = y_pos
        num_cols = len(data['table_headers'])
        col_width = CONTENT_WIDTH / num_cols

        # Header row
        header_height = Inches(0.5)
        for i, header in enumerate(data['table_headers']):
            header_box = slide.shapes.add_shape(
                1,  # Rectangle
                MARGIN + (i * col_width), table_top,
                col_width, header_height
            )
            header_box.fill.solid()
            header_box.fill.fore_color.rgb = PURPLE_CORE
            header_box.line.fill.background()  # No outline
            remove_shadow(header_box)  # Remove shadow completely

            tf = header_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = header
            p.font.name = 'Graphik'
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.LEFT

        # Content row
        content_row_height = Inches(0.9)
        y_content = table_top + header_height

        if data['table_content']:
            for i, cell_text in enumerate(data['table_content'][0]):
                cell_box = slide.shapes.add_textbox(
                    MARGIN + (i * col_width) + Inches(0.1), y_content + Inches(0.1),
                    col_width - Inches(0.2), content_row_height - Inches(0.2)
                )
                cf = cell_box.text_frame
                cf.word_wrap = True
                cf.vertical_anchor = MSO_ANCHOR.TOP
                cp = cf.paragraphs[0]
                cp.text = cell_text
                cp.font.name = 'Graphik'
                cp.font.size = Pt(11)
                cp.font.color.rgb = BLACK

        y_pos = y_content + content_row_height + Inches(0.35)

    # Bottom line callout
    if len(data['callouts']) > 2:
        callout = data['callouts'][2]
        bg_color = GRADIENT_GREEN_START if 'recommendation' in callout['type'] else GRADIENT_PURPLE_START
        border_color = GREEN if 'recommendation' in callout['type'] else PURPLE_CORE

        bg_box = slide.shapes.add_shape(
            1,  # Rectangle
            MARGIN, y_pos,
            CONTENT_WIDTH, Inches(0.6)
        )
        bg_box.fill.solid()
        bg_box.fill.fore_color.rgb = bg_color
        bg_box.line.fill.background()  # No outline
        remove_shadow(bg_box)  # Remove shadow completely

        title_width = Inches(2.0)
        title_box = slide.shapes.add_textbox(
            MARGIN + Inches(0.2), y_pos + Inches(0.12),
            title_width, Inches(0.35)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = callout['title']
        p.font.name = 'Graphik'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLACK

        content_box = slide.shapes.add_textbox(
            MARGIN + title_width + Inches(0.3), y_pos + Inches(0.08),
            CONTENT_WIDTH - title_width - Inches(0.5), Inches(0.44)
        )
        cf = content_box.text_frame
        cf.word_wrap = True
        cf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = cf.paragraphs[0]
        cp.text = callout['content']
        cp.font.name = 'Graphik'
        cp.font.size = Pt(11)
        cp.font.color.rgb = BLACK

        y_pos += Inches(0.75)

    # Footer with > symbol on left and info on right
    footer_y = prs.slide_height - Inches(0.7)  # Adjusted to prevent bleeding below slide

    # Greater than symbol (left)
    gt_box = slide.shapes.add_textbox(
        MARGIN, footer_y,
        Inches(0.5), Inches(0.4)
    )
    gt_frame = gt_box.text_frame
    p = gt_frame.paragraphs[0]
    p.text = '>'
    p.font.name = 'Graphik'
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = PURPLE_CORE

    # Footer info (right)
    footer_box = slide.shapes.add_textbox(
        prs.slide_width - Inches(4.5), footer_y,
        Inches(4.0), Inches(0.4)
    )
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = data['footer_contact']
    p.font.name = 'Graphik'
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT

    # Save
    prs.save(output_path)
    print(f"Created: {output_path}")


def main():
    """Process all 01-one-pager.html files"""

    tools = [
        ('Harness-Efficacy', 'Harness'),
        ('CoTester-Efficacy', 'CoTester'),
        ('Braintrust-Efficacy', 'Braintrust')
    ]

    for folder, tool_name in tools:
        html_path = Path(f'/home/christopher.g.roge/REPOS/00-TOOLS-RESEARCH/{folder}/01-one-pager.html')
        output_path = Path(f'/home/christopher.g.roge/REPOS/00-TOOLS-RESEARCH/{folder}/01-one-pager.pptx')

        if html_path.exists():
            print(f"\nProcessing {tool_name}...")
            data = parse_html_one_pager(html_path)
            create_pptx_slide(data, output_path, tool_name)
        else:
            print(f"Warning: {html_path} not found")

    print("\n✓ All PowerPoint decks created!")


if __name__ == '__main__':
    main()
